"""
Export av en simulering till olika dokumentformat:
  * "sida"    — enkel, självständig HTML-sida (öppnas i webbläsare)
  * "pptx"    — PowerPoint-presentation (python-pptx)
  * "rapport" — Word-dokument .docx (python-docx)

Alla format tar samma 'summary'-dict från frontend + valfri rapporttext (kan
AI-genereras på valt språk). En mottagare kan anges (dokumentet adresseras till den).
Returnerar (bytes, filnamn, mimetype).
"""

import html
import io
import json
from datetime import datetime


def _rows_arter(summary):
    return summary.get("arter", []) or []


def _param_items(summary):
    return list((summary.get("parametrar") or {}).items())


def _today():
    return datetime.now().strftime("%Y-%m-%d")


CONTACT_EMAIL = "hugo.wikstrom@gmail.com"

# Målgruppslägen för rapporten. Styr framing/ton (och AI-textens prompt via advisor).
MODES = {
    "nyfiken": {
        "namn": "Nyfiken",
        "lead": "Den här rapporten är skriven för dig som är nyfiken på Östersjön — "
                "enkelt och kort, utan fackspråk."},
    "generell": {
        "namn": "Generell",
        "lead": "En allmän, saklig rapport om körningens resultat."},
    "politik": {
        "namn": "Politik & beslut",
        "lead": "Ett beslutsunderlag för politik och förvaltning: strategier, argument, "
                "investeringar och samhällsvärden. Ekonomiska värden visas i både euro och "
                "lokal valuta per land."},
    "forskare": {
        "namn": "Forskning",
        "lead": "En fördjupning för forskare: mekanismer och samband, osäkerheter och "
                "kunskapsluckor, samt förslag på nästa steg i forskningen."},
}


def _mode(mode):
    return MODES.get(mode or "generell", MODES["generell"])


def _economy_rows(eco, mode):
    """
    Bygger ekonomirader (land, eur-sträng, lokal-sträng|None).
    I politik-läget läggs lokal valuta till per land.
    """
    politik = (mode == "politik")
    rows = []
    for land, v in eco.items():
        eur = f"{v} M€"
        loc = None
        if politik:
            try:
                from model import economics as E
                amt, ccy = E.local_amount(float(v), land)
                loc = "—" if ccy == "EUR" else f"{amt:,.0f} M{ccy}".replace(",", " ")
            except Exception:
                loc = None
        rows.append((str(land), eur, loc))
    return rows


# --- HaV-rapportens fasta avsnittstexter -------------------------------------
# Dispositionen följer Havs- och vattenmyndighetens rapportstruktur. Texterna är
# sakligt korrekta om modellen; den AI-genererade texten (report_text) läggs i
# Sammanfattningen. Nycklar kan översättas via 'strings' (i18n) men har svensk
# fallback här.
SECT = {
    "forord": (
        "Denna rapport är genererad av Eystrasalt — en öppen, pedagogisk digital "
        "tvilling av Östersjöns ekosystem. Rapportens disposition följer Havs- och "
        "vattenmyndighetens rapportstruktur. Eystrasalt är byggt av en enskild "
        "entusiast med stöd av AI och är varken en publikation från eller granskad "
        "av Havs- och vattenmyndigheten. Modellen är en litteraturförankrad "
        "förenkling avsedd för utforskning och lärande — inte en forskningsvaliderad "
        "prognosmodell."),
    "bakgrund": (
        "Östersjön är ett av världens mest övergödda och känsliga innanhav. "
        "Övergödning, klimatuppvärmning, syrefria bottnar och hårt fisketryck "
        "samverkar och har vid flera tillfällen utlöst regimskiften i näringsväven — "
        "mest känt torskens kollaps och storspiggens expansion. Att förstå hur "
        "dessa faktorer hänger ihop kräver att man ser hela kretsloppet, från "
        "näringssalter till toppredatorer."),
    "syfte": (
        "Syftet är att åskådliggöra hur Östersjöns ekosystem svarar på olika "
        "kombinationer av stress (klimat, övergödning, salthalt, fiske och jakt) "
        "och att jämföra förvaltningsstrategier under osäkerhet. Målet är ökad "
        "förståelse och en gemensam, litteraturförankrad utgångspunkt för samtal "
        "om havets framtid."),
    "avgransningar": (
        "Modellen är en förenklad box-modell med 17 kompartment i sex havsområden. "
        "Absoluta värden anges som modellindex i storleksordning g/m² (biomassa), "
        "mg/m³ (näring) och % mättnad (syre) — de är inte kalibrerade "
        "beståndsuppskattningar. Rumslig upplösning, individvariation, detaljerad "
        "hydrografi och ekonomiska data är kraftigt förenklade."),
    "modell": (
        "Ekosystemet beskrivs av kopplade differentialekvationer (en NPZD-modell "
        "utökad med fisk, kustrovfisk, lax, sjöfågel och säl) som löses numeriskt "
        "över vald tidshorisont. Biologiska hastigheter är temperaturberoende "
        "(Q10 = 2). Predation modelleras med mättande funktionellt svar (Holling "
        "typ II). Sex havsområden med salthaltsgradient utbyter näring och salt. "
        "Syret delas i yt- och bottenvatten för att fånga skiktning och syrefria "
        "bottnar. Kväve-sänkor (denitrifikation, utflöde, sedimentbegravning) "
        "sluter näringsbudgeten. Osäkra naturparametrar analyseras med en Monte "
        "Carlo-ensemble."),
    "diskussion": (
        "Resultaten ska tolkas kvalitativt — som riktningar och mekanismer snarare "
        "än exakta prognoser. Modellens styrka är att den visar hur delarna hänger "
        "ihop: hur övergödning och uppvärmning driver syrebrist, hur hårt torskfiske "
        "kan låsa fast ett regimskifte, och hur kustrovfisk och storspigg "
        "konkurrerar. Osäkerheten i naturparametrarna belyses i Monte Carlo-analysen "
        "och känslighetsanalysen (var mer forskning behövs)."),
}

SECT_TITLES = {
    "forord_h": "Förord", "sammanfattning_h": "Sammanfattning",
    "innehall_h": "Innehåll", "inledning_h": "1 Inledning",
    "bakgrund_h": "1.1 Bakgrund", "syfte_h": "1.2 Syfte och mål",
    "avgransningar_h": "1.3 Avgränsningar", "metod_h": "2 Metod",
    "modell_h": "2.1 Modellbeskrivning", "scenario_h": "2.2 Scenario och inställningar",
    "resultat_h": "3 Resultat", "res_halsa_h": "3.1 Ekosystemets hälsa",
    "res_arter_h": "3.2 Arter", "res_pyramid_h": "3.3 Näringspyramid",
    "res_uttak_h": "3.4 Uttag ur havet", "res_strategi_h": "3.5 Strategier och ekonomi",
    "diskussion_h": "4 Diskussion", "slutsatser_h": "5 Slutsatser",
    "referenser_h": "6 Referenser", "bilaga_h": "Bilaga A — Fullständiga inställningar",
}


def _S(strings, key):
    """Avsnittsrubrik: i18n-översättning om den finns, annars svensk fallback."""
    return strings.get(key, SECT_TITLES.get(key, key))


def _slutsatser(summary):
    """Punktvisa slutsatser härledda ur körningens data (ingen påhittad text)."""
    out = []
    hn = summary.get("halsa_nu")
    if hn is not None:
        out.append(f"Ekosystemets hälsoindex i slutet av körningen: {hn} / 100.")
    mc = summary.get("mc")
    if mc and mc.get("basta_namn"):
        out.append(f"Strategin med bäst återhämtning i Monte Carlo-analysen: "
                   f"{mc.get('basta_namn')}.")
    out.append("Resultaten är kvalitativa riktningar från en förenklad modell, "
               "inte prognoser — se Avgränsningar och Diskussion.")
    return out


def _references():
    """Referenslista: modellens litteraturbibliotek + ev. användarens rapporter."""
    refs = []
    try:
        from ai.seed_reports import LIBRARY
        refs.extend(titel for titel, _ in LIBRARY)
    except Exception:
        pass
    try:
        from ai import reports
        for r in reports.list_reports():
            t = r.get("titel", "")
            if t and t not in refs:
                refs.append(t)
    except Exception:
        pass
    return refs


# --- Lägesspecifikt innehåll (så lägena skiljer sig ÄVEN utan AI-text) -------
def _halsa_ord(hn):
    try:
        hn = float(hn)
    except (TypeError, ValueError):
        return "okänt"
    if hn >= 75:
        return "Havet mår förhållandevis bra i den här körningen"
    if hn >= 55:
        return "Havet är påverkat men fungerar"
    if hn >= 35:
        return "Havet visar tydliga tecken på stress"
    return "Havet är hårt pressat"


def _top_movers(arter, n=3):
    """Plain-språk: de arter som ändrades mest (för nyfiken-läget)."""
    scored = []
    for a in arter:
        try:
            s, e = float(a.get("start")), float(a.get("slut"))
        except (TypeError, ValueError):
            continue
        scored.append((abs(e - s), str(a.get("namn", "")), s, e))
    scored.sort(reverse=True)
    out = []
    for _, namn, s, e in scored[:n]:
        rikt = "ökade" if e > s else ("minskade" if e < s else "var stabil")
        out.append(f"{namn} {rikt} (från {s:g} till {e:g}).")
    return out


def _nyfiken_fallback(summary):
    hn = summary.get("halsa_nu")
    if hn is not None:
        return (f"I den här körningen hamnade havets hälsa på {hn} av 100. "
                "Här är det som förändrades mest bland arterna.")
    return "Här är en enkel överblick av din körning."


def _politik_fallback(summary):
    hn = summary.get("halsa_nu")
    mc = summary.get("mc") or {}
    bits = []
    if hn is not None:
        bits.append(f"Ekosystemets hälsa i körningen: {hn}/100.")
    if mc.get("basta_namn"):
        bits.append(f"Strategin med bäst återhämtning: {mc['basta_namn']}.")
    eco = mc.get("ekonomi_per_land") or {}
    if eco:
        try:
            tot = sum(float(v) for v in eco.values())
            bits.append(f"Sammanlagt beräknat årligt värde av åtgärden: cirka {tot:g} M€/år, "
                        "fördelat på länderna enligt tabellen nedan.")
        except (TypeError, ValueError):
            pass
    bits.append("Nedan följer handlingsalternativ och argument för beslut.")
    return " ".join(bits)


def _forskare_fallback(summary):
    hn = summary.get("halsa_nu")
    s = "Sammanfattning av körningen för vidare analys. "
    if hn is not None:
        s += f"Hälsoindex: {hn}/100. "
    s += ("Resultaten är kvalitativa riktningar från en förenklad, litteraturförankrad "
          "modell — fokusera på mekanismer, osäkerheter och känslighet nedan.")
    return s


POLITIK_LEVERS = [
    "Minska närsaltbelastningen (kväve och fosfor) enligt HELCOM:s aktionsplan (BSAP) — "
    "störst långsiktig effekt, men havet svarar trögt på grund av intern belastning.",
    "Anpassa fiskeuttaget efter beståndens status (t.ex. torskstopp vid kollaps) och "
    "samordna kvoter mellan länderna.",
    "Åtgärder i kustzonen mot spigg-dominans: skydda lek- och uppväxtområden för abborre och gädda.",
    "Investera i övervakning och åtgärder samordnat via HELCOM — havet delas av nio länder "
    "och vinsterna fördelas ojämnt (se tabellen).",
    "Argument: kostnaden för att inte agera (förlorade ekosystemtjänster, fiske och turism) "
    "överstiger på sikt åtgärdskostnaden.",
]
FORSKARE_GAPS = [
    "Intern fosforbelastning från syrefria bottnar — styrkan i den självförstärkande onda cirkeln.",
    "Storspiggens larvpredation på torsk, abborre och gädda (wasp-waist-styrkan).",
    "Frekvens och effekt av större saltvatteninflöden (MBI) samt skiktningens klimatberoende.",
    "Kalibrering av modellens relativa biomassaenheter mot uppmätta beståndsuppskattningar.",
]
FORSKARE_NEXT = [
    "Prioritera mätserier av bottensyre och fosforflöden i de djupa bassängerna.",
    "Fältstudier av spiggens predation på fiskyngel längs salthaltsgradienten.",
    "Koppla modellen till hydrografiska data (inflöden, temperatur, skiktning) för bättre syredynamik.",
    "Validera de ekonomiska koefficienterna mot BalticSTERN/HELCOM för robustare värdering.",
]


# --- Enkel HTML-sida ---------------------------------------------------------
def _build_html(summary, report_text, recipient, title, strings, mode="generell"):
    esc = html.escape
    arter = _rows_arter(summary)
    params = _param_items(summary)
    trofi = summary.get("trofi") or {}
    niv = (trofi.get("nivaer_slut") or {})
    uttak = summary.get("uttak_slut") or {}
    mc = summary.get("mc")
    md = _mode(mode)

    S = lambda k: _S(strings, k)
    parts = [f"""<!doctype html><html lang="sv"><head><meta charset="utf-8">
<meta name="viewport" content="width=device-width, initial-scale=1">
<title>{esc(title)}</title>
<style>
 :root{{--ink:#12293f;--sea:#0e7490;--sea2:#155e75;--line:#cfe3ef;--soft:#eef6fb;--muted:#5b7791}}
 *{{box-sizing:border-box}}
 body{{font-family:'Georgia','Iowan Old Style',ui-serif,serif;max-width:760px;margin:0 auto;
   padding:40px 26px 64px;color:var(--ink);background:#fff;line-height:1.7;font-size:17px}}
 h1,h2,h3{{font-family:system-ui,'Segoe UI',Roboto,sans-serif;line-height:1.25}}
 h1{{color:var(--sea);font-size:2rem;margin:.2em 0}}
 h2{{color:var(--sea);border-bottom:2px solid var(--line);padding-bottom:6px;margin-top:2.2em;font-size:1.4rem}}
 h3{{color:var(--sea2);margin-top:1.5em;font-size:1.12rem}}
 p{{margin:.7em 0}}
 table{{border-collapse:collapse;width:100%;margin:14px 0;font-family:system-ui,sans-serif;font-size:.95rem}}
 th,td{{border:1px solid var(--line);padding:8px 11px;text-align:right}}
 th:first-child,td:first-child{{text-align:left}}
 thead th{{background:var(--soft);color:var(--sea2)}}
 tbody tr:nth-child(even){{background:#fafdff}}
 .muted{{color:var(--muted)}} .badge{{background:var(--soft);border-radius:8px;padding:3px 10px;font-weight:600}}
 .to{{background:var(--soft);border-left:4px solid var(--sea);border-radius:6px;padding:10px 14px;margin:14px 0}}
 .toc{{background:var(--soft);border:1px solid var(--line);border-radius:10px;padding:12px 22px;
   font-family:system-ui,sans-serif;font-size:.95rem}}
 .toc a{{color:var(--sea);text-decoration:none}} .toc a:hover{{text-decoration:underline}}
 .cover{{text-align:center;padding:20px 0 8px;border-bottom:1px solid var(--line);margin-bottom:8px}}
 .modebadge{{display:inline-block;background:var(--sea);color:#fff;border-radius:999px;
   padding:4px 14px;font-family:system-ui,sans-serif;font-size:.85rem;font-weight:600;margin-top:6px}}
 .lead{{font-size:1.05rem;color:var(--sea2);background:var(--soft);border-radius:10px;
   padding:14px 18px;margin:16px 0}}
 .summary-card{{background:var(--soft);border:1px solid var(--line);border-radius:12px;padding:6px 22px;margin:14px 0}}
 @media print{{body{{padding:0;font-size:12pt}} h2{{page-break-after:avoid}} table{{page-break-inside:avoid}}}}
</style></head><body>
<div class="cover">
<h1>🌊 {esc(title)}</h1>
<p class="muted">{esc(strings.get("undertitel","Simulering av Östersjöns ekosystem — digital tvilling"))}</p>
<div class="modebadge">{esc(strings.get("mode_"+ (mode or "generell"), md["namn"]))}</div>
<p class="muted">Rapportstruktur enligt Havs- och vattenmyndighetens mall · Datum: {esc(_today())}</p>
</div>
<p class="lead">{esc(md["lead"])}</p>"""]
    if recipient:
        parts.append(f'<div class="to"><b>{esc(strings.get("till","Till"))}:</b> {esc(recipient)}</div>')

    hn = summary.get("halsa_nu")

    # ---- Återanvändbara byggblock (samma data, olika disposition per läge) ---
    def summary_block(fallback):
        out = ['<div class="summary-card">']
        if report_text:
            for para in report_text.split("\n"):
                para = para.strip()
                if not para:
                    continue
                if para.startswith("#"):
                    out.append(f"<h3>{esc(para.lstrip('# ').strip())}</h3>")
                elif para.startswith("**") and para.endswith("**") and len(para) > 4:
                    out.append(f"<h3>{esc(para.strip('*').strip())}</h3>")
                else:
                    out.append(f"<p>{esc(para)}</p>")
        else:
            out.append(f"<p>{esc(fallback)}</p>")
        out.append("</div>")
        return out

    def params_table():
        if not params:
            return ["<p class='muted'>Standardinställningar (baslinje).</p>"]
        out = ['<table><thead><tr><th>Parameter</th><th>Värde</th></tr></thead><tbody>']
        for k, v in params:
            out.append(f"<tr><td>{esc(str(k))}</td><td>{esc(str(v))}</td></tr>")
        out.append("</tbody></table>")
        return out

    def arter_table():
        if not arter:
            return []
        out = ['<table><thead><tr><th>Art</th><th>Start</th><th>Slut</th>'
               '<th>Enhet</th></tr></thead><tbody>']
        for a in arter:
            out.append(f"<tr><td>{esc(str(a.get('namn','')))}</td><td>{esc(str(a.get('start','')))}</td>"
                       f"<td>{esc(str(a.get('slut','')))}</td><td class='muted'>{esc(str(a.get('enhet','')))}</td></tr>")
        out.append("</tbody></table>")
        return out

    def pyramid_table():
        if not niv:
            return []
        out = ['<table><thead><tr><th>Nivå</th><th>g/m²</th></tr></thead><tbody>']
        for namn, val in niv.items():
            out.append(f"<tr><td>{esc(str(namn))}</td><td>{esc(str(val))}</td></tr>")
        out.append("</tbody></table>")
        return out

    def uttak_table():
        if not uttak:
            return []
        out = ['<table><thead><tr><th>Uttag</th><th>g/m²/år</th></tr></thead><tbody>']
        labels = {"fiske": "Fiske (människa)", "sal": "Säl", "skarv": "Skarv/sjöfågel",
                  "atervinning": "Återvinning (nedbrytning→näring)"}
        for k, lab in labels.items():
            if k in uttak:
                out.append(f"<tr><td>{esc(lab)}</td><td>{esc(str(uttak[k]))}</td></tr>")
        out.append("</tbody></table>")
        return out

    def economy_block(politik=False):
        if not mc:
            return []
        out = [f"<p><b>{esc(str(mc.get('basta_namn','')))}</b></p>"]
        eco = mc.get("ekonomi_per_land") or {}
        if eco:
            rows = _economy_rows(eco, mode)
            if politik:
                out.append("<table><thead><tr><th>Land</th><th>Euro/år</th>"
                           "<th>Lokal valuta/år</th></tr></thead><tbody>")
                for land, eur, loc in rows:
                    out.append(f"<tr><td>{esc(land)}</td><td>{esc(eur)}</td>"
                               f"<td>{esc(loc or '—')}</td></tr>")
                out.append("</tbody></table>")
                out.append('<p class="muted">Lokal valuta omräknad från euro med ungefärliga '
                           'växelkurser — storleksordningar för jämförelse.</p>')
            else:
                out.append("<table><thead><tr><th>Land</th><th>M€/år</th></tr></thead><tbody>")
                for land, eur, _loc in rows:
                    out.append(f"<tr><td>{esc(land)}</td><td>{esc(eur)}</td></tr>")
                out.append("</tbody></table>")
        return out

    def health_badge():
        if hn is None:
            return []
        return [f'<p><span class="badge">{esc(str(hn))} / 100</span> (index 0–100)</p>']

    def ul(items):
        return ["<ul>"] + [f"<li>{esc(x)}</li>" for x in items] + ["</ul>"]

    def ol(items):
        return ["<ol>"] + [f"<li>{esc(x)}</li>" for x in items] + ["</ol>"]

    # ---- Lägesspecifik disposition ------------------------------------------
    if mode == "nyfiken":
        parts.append("<h2>Så mår havet</h2>")
        if hn is not None:
            parts.append(f'<p style="font-size:2.6rem;margin:.1em 0"><b>{esc(str(hn))}</b>'
                         f'<span class="muted" style="font-size:1rem"> / 100</span></p>')
            parts.append(f"<p>{esc(_halsa_ord(hn))}.</p>")
        parts.append("<h2>Vad hände i din körning?</h2>")
        parts += summary_block(_nyfiken_fallback(summary))
        movers = _top_movers(arter)
        if movers:
            parts += ul(movers)
        parts.append("<h2>Varför spelar det roll?</h2>")
        parts.append("<p>Östersjön är ett känsligt innanhav. När en del ändras — övergödning, "
                     "värme eller fiske — påverkas hela näringsväven, ända upp till torsk, fågel "
                     "och säl. Ett friskare hav ger mer fisk, klarare vatten och ett rikare liv "
                     "för oss alla runt kusten.</p>")
        parts.append('<p class="muted">Vill du gräva djupare? Byt läge till „Generell”, '
                     '„Politik” eller „Forskare” när du exporterar.</p>')

    elif mode == "politik":
        parts.append("<h2>Sammanfattning för beslutsfattare</h2>")
        parts += summary_block(_politik_fallback(summary))
        parts.append("<h2>Ekonomiskt värde per land</h2>")
        if mc:
            parts += economy_block(politik=True)
        else:
            parts.append('<p class="muted">Kör en Monte Carlo-jämförelse för ekonomiska '
                         'värden per land.</p>')
        parts.append("<h2>Strategier och handlingsalternativ</h2>")
        parts += ul(POLITIK_LEVERS)
        parts.append("<h2>Resultat i korthet</h2>")
        parts += health_badge()
        parts += arter_table()
        parts.append("<h2>Slutsatser</h2>")
        parts += ul(_slutsatser(summary))
        parts.append(f'<h2>Om underlaget</h2><p>{esc(SECT["avgransningar"])}</p>')

    elif mode == "forskare":
        parts.append("<h2>Sammanfattning</h2>")
        parts += summary_block(_forskare_fallback(summary))
        parts.append(f'<h2>Metod och modell</h2><p>{esc(SECT["modell"])}</p>')
        parts.append("<h3>Scenario och inställningar</h3>")
        parts += params_table()
        parts.append("<h2>Resultat</h2>")
        parts += health_badge()
        parts += arter_table()
        if niv:
            parts.append("<h3>Näringspyramid (g/m²)</h3>")
            parts += pyramid_table()
        if uttak:
            parts.append("<h3>Uttag ur havet</h3>")
            parts += uttak_table()
        if mc:
            parts.append("<h3>Strategi och ekonomi</h3>")
            parts += economy_block()
        parts.append(f'<h2>Osäkerheter och känslighet</h2><p>{esc(SECT["avgransningar"])}</p>')
        parts.append("<p>Osäkra naturparametrar analyseras i appens Monte Carlo- och "
                     "känslighetsanalys. De största osäkerheterna rör:</p>")
        parts += ul(FORSKARE_GAPS)
        parts.append("<h2>Kunskapsluckor och nästa steg</h2>")
        parts += ol(FORSKARE_NEXT)
        parts.append(f'<h2>Referenser</h2>')
        parts += ol(_references())

    else:  # generell — fullständig HaV-struktur
        parts.append(f'<h2 id="forord">{esc(S("forord_h"))}</h2><p>{esc(SECT["forord"])}</p>')
        parts.append(f'<h2 id="sammanfattning">{esc(S("sammanfattning_h"))}</h2>')
        parts += summary_block("(Ingen sammanfattande text genererades — kryssa i "
                               "„Låt AI:n skriva rapporttext” och kör igen.)")
        parts.append(f'<h2>{esc(S("innehall_h"))}</h2><div class="toc">'
                     f'<a href="#forord">{esc(S("forord_h"))}</a><br>'
                     f'<a href="#sammanfattning">{esc(S("sammanfattning_h"))}</a><br>'
                     f'<a href="#inledning">{esc(S("inledning_h"))}</a><br>'
                     f'<a href="#metod">{esc(S("metod_h"))}</a><br>'
                     f'<a href="#resultat">{esc(S("resultat_h"))}</a><br>'
                     f'<a href="#diskussion">{esc(S("diskussion_h"))}</a><br>'
                     f'<a href="#slutsatser">{esc(S("slutsatser_h"))}</a><br>'
                     f'<a href="#referenser">{esc(S("referenser_h"))}</a></div>')
        parts.append(f'<h2 id="inledning">{esc(S("inledning_h"))}</h2>')
        parts.append(f'<h3>{esc(S("bakgrund_h"))}</h3><p>{esc(SECT["bakgrund"])}</p>')
        parts.append(f'<h3>{esc(S("syfte_h"))}</h3><p>{esc(SECT["syfte"])}</p>')
        parts.append(f'<h3>{esc(S("avgransningar_h"))}</h3><p>{esc(SECT["avgransningar"])}</p>')
        parts.append(f'<h2 id="metod">{esc(S("metod_h"))}</h2>')
        parts.append(f'<h3>{esc(S("modell_h"))}</h3><p>{esc(SECT["modell"])}</p>')
        parts.append(f'<h3>{esc(S("scenario_h"))}</h3>')
        parts += params_table()
        parts.append(f'<h2 id="resultat">{esc(S("resultat_h"))}</h2>')
        if hn is not None:
            parts.append(f'<h3>{esc(S("res_halsa_h"))}</h3>')
            parts += health_badge()
        if arter:
            parts.append(f'<h3>{esc(S("res_arter_h"))}</h3>')
            parts += arter_table()
        if niv:
            parts.append(f'<h3>{esc(S("res_pyramid_h"))}</h3>')
            parts += pyramid_table()
        if uttak:
            parts.append(f'<h3>{esc(S("res_uttak_h"))}</h3>')
            parts += uttak_table()
        if mc:
            parts.append(f'<h3>{esc(S("res_strategi_h"))}</h3>')
            parts += economy_block(politik=False)
        parts.append(f'<h2 id="diskussion">{esc(S("diskussion_h"))}</h2><p>{esc(SECT["diskussion"])}</p>')
        parts.append(f'<h2 id="slutsatser">{esc(S("slutsatser_h"))}</h2>')
        parts += ul(_slutsatser(summary))
        parts.append(f'<h2 id="referenser">{esc(S("referenser_h"))}</h2>')
        parts += ol(_references())

    parts.append(f'<p class="muted" style="margin-top:36px;border-top:1px solid var(--line);'
                 f'padding-top:12px">Genererad av Eystrasalt — open source digital tvilling '
                 f'för Östersjön. {CONTACT_EMAIL}</p>')
    parts.append("</body></html>")
    return "\n".join(parts).encode("utf-8")


# --- PowerPoint --------------------------------------------------------------
def _build_pptx(summary, report_text, recipient, title, strings, mode="generell"):
    from pptx import Presentation
    from pptx.util import Inches, Pt

    prs = Presentation()
    blank = prs.slide_layouts[6]
    title_layout = prs.slide_layouts[0]

    # Titelbild
    s = prs.slides.add_slide(title_layout)
    s.shapes.title.text = title
    md = _mode(mode)
    sub = strings.get("undertitel", "Simulering av Östersjöns ekosystem")
    sub += f"\n{md['namn']}"
    if recipient:
        sub += f"\n{strings.get('till','Till')}: {recipient}"
    s.placeholders[1].text = sub

    def bullet_slide(rubrik, rader):
        sl = prs.slides.add_slide(blank)
        tb = sl.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(9), Inches(1))
        tf = tb.text_frame; tf.text = rubrik
        tf.paragraphs[0].font.size = Pt(28); tf.paragraphs[0].font.bold = True
        body = sl.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(8.8), Inches(5))
        bf = body.text_frame; bf.word_wrap = True
        for i, r in enumerate(rader):
            p = bf.paragraphs[0] if i == 0 else bf.add_paragraph()
            p.text = "• " + r; p.font.size = Pt(16)

    hn = summary.get("halsa_nu")
    if hn is not None:
        bullet_slide(strings.get("halsa", "Ekosystemets hälsa"),
                     [f"Hälsoindex: {hn} / 100"])
    params = _param_items(summary)
    if params:
        bullet_slide(strings.get("installningar", "Inställningar"),
                     [f"{k}: {v}" for k, v in params])
    arter = _rows_arter(summary)
    if arter:
        bullet_slide(strings.get("arter", "Arter (start → slut)"),
                     [f"{a.get('namn','')}: {a.get('start','')} → {a.get('slut','')} {a.get('enhet','')}"
                      for a in arter])
    niv = (summary.get("trofi") or {}).get("nivaer_slut") or {}
    if niv:
        bullet_slide(strings.get("pyramid", "Näringspyramid"),
                     [f"{namn}: {val}" for namn, val in niv.items()])
    mc = summary.get("mc")
    if mc:
        rows = [f"Bästa strategi: {mc.get('basta_namn','')}"]
        for land, eur, loc in _economy_rows(mc.get("ekonomi_per_land") or {}, mode):
            rows.append(f"{land}: {eur}/år" + (f"  ({loc}/år)" if loc and loc != "—" else ""))
        bullet_slide(strings.get("strategi", "Strategi & ekonomi"), rows)
    if report_text:
        for chunk_i, chunk in enumerate(_chunks(report_text, 6)):
            bullet_slide(strings.get("sammanfattning", "Sammanfattning")
                         + ("" if chunk_i == 0 else " (forts.)"), chunk)

    buf = io.BytesIO(); prs.save(buf); return buf.getvalue()


def _chunks(text, n):
    paras = [p.strip() for p in text.split("\n") if p.strip()]
    for i in range(0, len(paras), n):
        yield paras[i:i + n]


# --- Word-rapport (struktur enligt Havs- och vattenmyndighetens rapportmall) --
def _toc_field(doc):
    """Infogar ett innehållsförteckningsfält som Word fyller i när dokumentet öppnas."""
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement
    p = doc.add_paragraph()
    run = p.add_run()
    fld = OxmlElement("w:fldSimple")
    fld.set(qn("w:instr"), r'TOC \o "1-2" \h \z \u')
    hint = OxmlElement("w:r"); ht = OxmlElement("w:t")
    ht.text = "Högerklicka och välj „Uppdatera fält” för att bygga innehållsförteckningen."
    hint.append(ht); fld.append(hint)
    run._r.addnext(fld)


def _build_docx(summary, report_text, recipient, title, strings, mode="generell"):
    from docx import Document
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.shared import Pt, RGBColor

    doc = Document()
    S = lambda k: _S(strings, k)
    md = _mode(mode)

    # --- Titelsida -----------------------------------------------------------
    for _ in range(6):
        doc.add_paragraph()
    h = doc.add_heading(title, 0); h.alignment = WD_ALIGN_PARAGRAPH.CENTER
    sub = doc.add_paragraph(strings.get("undertitel",
          "Simulering av Östersjöns ekosystem — digital tvilling"))
    sub.alignment = WD_ALIGN_PARAGRAPH.CENTER
    meta = doc.add_paragraph(); meta.alignment = WD_ALIGN_PARAGRAPH.CENTER
    meta.add_run(f"Läge: {md['namn']}\n"
                 f"Rapportstruktur enligt Havs- och vattenmyndighetens mall\n"
                 f"Datum: {_today()}\n"
                 + (f"Mottagare: {recipient}\n" if recipient else "")
                 + "Utgivare: Eystrasalt — öppen digital tvilling för Östersjön")
    doc.add_page_break()

    # --- Förord --------------------------------------------------------------
    doc.add_heading(S("forord_h"), 1)
    lead = doc.add_paragraph(); lead.add_run(md["lead"]).italic = True
    doc.add_paragraph(SECT["forord"])

    # --- Sammanfattning (AI-text) -------------------------------------------
    doc.add_heading(S("sammanfattning_h"), 1)
    if report_text:
        for para in report_text.split("\n"):
            para = para.strip()
            if para:
                doc.add_paragraph(para)
    else:
        doc.add_paragraph("(Ingen sammanfattande text genererades för denna körning.)")

    # --- Innehåll ------------------------------------------------------------
    doc.add_heading(S("innehall_h"), 1)
    _toc_field(doc)
    doc.add_page_break()

    # --- 1 Inledning ---------------------------------------------------------
    doc.add_heading(S("inledning_h"), 1)
    doc.add_heading(S("bakgrund_h"), 2); doc.add_paragraph(SECT["bakgrund"])
    doc.add_heading(S("syfte_h"), 2); doc.add_paragraph(SECT["syfte"])
    doc.add_heading(S("avgransningar_h"), 2); doc.add_paragraph(SECT["avgransningar"])

    # --- 2 Metod -------------------------------------------------------------
    doc.add_heading(S("metod_h"), 1)
    doc.add_heading(S("modell_h"), 2); doc.add_paragraph(SECT["modell"])
    doc.add_heading(S("scenario_h"), 2)
    params = _param_items(summary)
    if params:
        t = doc.add_table(rows=1, cols=2); t.style = "Light Grid Accent 1"
        t.rows[0].cells[0].text, t.rows[0].cells[1].text = "Parameter", "Värde"
        for k, v in params:
            c = t.add_row().cells; c[0].text = str(k); c[1].text = str(v)
    else:
        doc.add_paragraph("Standardinställningar (baslinje).")

    # --- 3 Resultat ----------------------------------------------------------
    doc.add_heading(S("resultat_h"), 1)
    hn = summary.get("halsa_nu")
    if hn is not None:
        doc.add_heading(S("res_halsa_h"), 2)
        doc.add_paragraph(f"Hälsoindex (0–100) i slutet av körningen: {hn} / 100.")
    arter = _rows_arter(summary)
    if arter:
        doc.add_heading(S("res_arter_h"), 2)
        t = doc.add_table(rows=1, cols=4); t.style = "Light Grid Accent 1"
        hdr = t.rows[0].cells
        hdr[0].text, hdr[1].text, hdr[2].text, hdr[3].text = "Art", "Start", "Slut", "Enhet"
        for a in arter:
            c = t.add_row().cells
            c[0].text = str(a.get("namn", "")); c[1].text = str(a.get("start", ""))
            c[2].text = str(a.get("slut", "")); c[3].text = str(a.get("enhet", ""))
    niv = (summary.get("trofi") or {}).get("nivaer_slut") or {}
    if niv:
        doc.add_heading(S("res_pyramid_h"), 2)
        for namn, val in niv.items():
            doc.add_paragraph(f"{namn}: {val} g/m²", style="List Bullet")
    uttak = summary.get("uttak_slut") or {}
    if uttak:
        doc.add_heading(S("res_uttak_h"), 2)
        labels = {"fiske": "Fiske (människa)", "sal": "Säl", "skarv": "Skarv/sjöfågel",
                  "atervinning": "Återvinning (nedbrytning→näring)"}
        for k, lab in labels.items():
            if k in uttak:
                doc.add_paragraph(f"{lab}: {uttak[k]} g/m²/år", style="List Bullet")
    mc = summary.get("mc")
    if mc:
        doc.add_heading(S("res_strategi_h"), 2)
        doc.add_paragraph(f"Strategi med bäst återhämtning: {mc.get('basta_namn','')}")
        eco = mc.get("ekonomi_per_land") or {}
        if mode == "politik" and eco:
            t = doc.add_table(rows=1, cols=3); t.style = "Light Grid Accent 1"
            h0 = t.rows[0].cells
            h0[0].text, h0[1].text, h0[2].text = "Land", "Euro/år", "Lokal valuta/år"
            for land, eur, loc in _economy_rows(eco, mode):
                c = t.add_row().cells
                c[0].text, c[1].text, c[2].text = land, eur, (loc or "—")
            doc.add_paragraph("Lokal valuta omräknad från euro med ungefärliga växelkurser "
                              "(storleksordningar för jämförelse).").runs[0].italic = True
        else:
            for land, eur, _loc in _economy_rows(eco, mode):
                doc.add_paragraph(f"{land}: {eur}/år", style="List Bullet")

    # --- 4 Diskussion --------------------------------------------------------
    doc.add_heading(S("diskussion_h"), 1)
    doc.add_paragraph(SECT["diskussion"])

    # --- 5 Slutsatser --------------------------------------------------------
    doc.add_heading(S("slutsatser_h"), 1)
    for s in _slutsatser(summary):
        doc.add_paragraph(s, style="List Bullet")

    # --- 6 Referenser --------------------------------------------------------
    doc.add_heading(S("referenser_h"), 1)
    for ref in _references():
        doc.add_paragraph(ref, style="List Bullet")

    # --- Bilaga A ------------------------------------------------------------
    if params:
        doc.add_heading(S("bilaga_h"), 1)
        t = doc.add_table(rows=1, cols=2); t.style = "Light Grid Accent 1"
        t.rows[0].cells[0].text, t.rows[0].cells[1].text = "Parameter", "Värde"
        for k, v in params:
            c = t.add_row().cells; c[0].text = str(k); c[1].text = str(v)

    doc.add_paragraph()
    p = doc.add_paragraph("Genererad av Eystrasalt — open source digital tvilling "
                          f"för Östersjön. {CONTACT_EMAIL}")
    p.runs[0].font.size = Pt(9)

    buf = io.BytesIO(); doc.save(buf); return buf.getvalue()


# --- Publikt API -------------------------------------------------------------
def export(fmt, summary, report_text=None, recipient="", title="Eystrasalt — Östersjörapport",
           strings=None, mode="generell"):
    """Bygger dokumentet. Returnerar (bytes, filnamn, mimetype)."""
    strings = strings or {}
    summary = summary or {}
    if fmt == "sida":
        data = _build_html(summary, report_text, recipient, title, strings, mode)
        return data, "eystrasalt.html", "text/html"
    if fmt == "pptx":
        data = _build_pptx(summary, report_text, recipient, title, strings, mode)
        return data, "eystrasalt.pptx", \
            "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    if fmt == "rapport":
        data = _build_docx(summary, report_text, recipient, title, strings, mode)
        return data, "eystrasalt.docx", \
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
    raise ValueError(f"Okänt format: {fmt}")
